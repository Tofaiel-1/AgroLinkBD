import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation_deck(output_filename):
    prs = Presentation()
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Color Palette
    PRIMARY_EMERALD = RGBColor(27, 94, 32)      # #1B5E20
    DARK_BG = RGBColor(18, 30, 22)              # #121E16 Dark Forest
    DEEP_TEAL = RGBColor(0, 121, 107)           # #00796B
    LIGHT_BG = RGBColor(245, 247, 245)          # #F5F7F5
    CARD_BG = RGBColor(255, 255, 255)           # #FFFFFF
    CARD_BORDER = RGBColor(200, 225, 204)       # #C8E1CC
    CARD_DARK_BG = RGBColor(28, 45, 33)         # #1C2D21
    GOLD_ACCENT = RGBColor(245, 124, 0)         # #F57C00
    TEXT_DARK = RGBColor(33, 33, 33)            # #212121
    TEXT_MUTED = RGBColor(117, 117, 117)        # #757575
    TEXT_LIGHT = RGBColor(240, 240, 240)        # #F0F0F0
    WHITE = RGBColor(255, 255, 255)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="INNOVATION FOR BANGLADESH 2026 | AGROLINKBD"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = GOLD_ACCENT

        # Header title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.7), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = PRIMARY_EMERALD

    def add_card(slide, left, top, width, height, title=None, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.2)
        else:
            shape.line.fill.background()
            
        if title:
            # Add card title header inside
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_EMERALD
        return shape

    # =========================================================================
    # SLIDE 1: TITLE & HERO SLIDE (DARK THEME)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, DARK_BG)

    # Hero Badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.0), Inches(4.5), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(38, 70, 48)
    badge.line.color.rgb = GOLD_ACCENT
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "🏆 NATIONAL INNOVATION FOR BANGLADESH 2026"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = GOLD_ACCENT
    p_b.alignment = PP_ALIGN.CENTER

    # Title
    tbox = s1.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.6))
    tf = tbox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "AgroLinkBD"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "এগ্রোলিংক বিডি: সমন্বিত স্মার্ট কৃষি ও মৎস্য বাণিজ্য সুপার-ইকোসিস্টেম"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(129, 199, 132) # Soft green

    # Subtitle
    sub_box = s1.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.9))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "A Full-Stack Generative AI, Agri-Fintech & Decentralized B2B Wholesale Trading Platform\nEmpowering 7 Interconnected Stakeholders across Bangladesh."
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = TEXT_LIGHT

    # Metric Highlight Cards (3 Bottom Cards)
    m1 = add_card(s1, Inches(0.8), Inches(4.7), Inches(3.6), Inches(1.8), bg_color=CARD_DARK_BG, border_color=PRIMARY_EMERALD)
    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(3.2), Inches(1.4))
    tf1 = tb1.text_frame
    p = tf1.paragraphs[0]
    p.text = "+25% to 40%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD_ACCENT
    p = tf1.add_paragraph()
    p.text = "Net Farm-Gate Income Increase via Live Auctions & Corporate Buyback Contracts"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_LIGHT

    m2 = add_card(s1, Inches(4.8), Inches(4.7), Inches(3.6), Inches(1.8), bg_color=CARD_DARK_BG, border_color=PRIMARY_EMERALD)
    tb2 = s1.shapes.add_textbox(Inches(5.0), Inches(4.9), Inches(3.2), Inches(1.4))
    tf2 = tb2.text_frame
    p = tf2.paragraphs[0]
    p.text = "-30% Waste"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(76, 175, 80) # Vibrant Green
    p = tf2.add_paragraph()
    p.text = "Reduction in Aquatic & Crop Disease Loss & Perishable Spoilage via Gemini Vision AI"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_LIGHT

    m3 = add_card(s1, Inches(8.8), Inches(4.7), Inches(3.7), Inches(1.8), bg_color=CARD_DARK_BG, border_color=PRIMARY_EMERALD)
    tb3 = s1.shapes.add_textbox(Inches(9.0), Inches(4.9), Inches(3.3), Inches(1.4))
    tf3 = tb3.text_frame
    p = tf3.paragraphs[0]
    p.text = "100% Inclusion"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(64, 196, 255) # Light Blue
    p = tf3.add_paragraph()
    p.text = "Unbanked Farmers Empowered via NFC Smart AgroCards & Automated Bank Loan KYC Files"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_LIGHT

    # Presenter tag
    pres_box = s1.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    p_pres = pres_box.text_frame.paragraphs[0]
    p_pres.text = "University Innovation Team | Track: Smart Agriculture, AI & Fintech | National Finalist (Top 150)"
    p_pres.font.size = Pt(11)
    p_pres.font.italic = True
    p_pres.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: THE MACRO CRISIS IN BANGLADESH (PROBLEM STATEMENT)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, LIGHT_BG)
    add_header(s2, "The Burning Crisis in Bangladesh Agriculture & Fisheries")

    cards_data_s2 = [
        ("1. Multi-Tier Middleman Exploitation", 
         "• Farmers lose 45%–60% profit to Faria, Paikar & Aratdars\n• Extreme price asymmetry: Farmer gets ৳15/kg, Dhaka consumer pays ৳65/kg\n• Absence of transparent open bidding floors forces distress selling.",
         Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.5)),
        ("2. Devastating Disease & Chemical Losses",
         "• ৳4,000+ Crore annual loss in fish & crop mortality\n• Over-the-counter quack chemical poisoning ruins pond ecosystems\n• Absence of instant pathology causes entire harvest wipeouts in 48 hrs.",
         Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.5)),
        ("3. Severe Post-Harvest Perishable Waste",
         "• 25%–30% of harvested vegetables & fish perish in transit\n• Fragmented rural upazila logistics with no cold-chain routing\n• Smallholders unable to rent high-tonnage transport independently.",
         Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6)),
        ("4. Rural Financial Exclusion & 'Dadon' Trap",
         "• 85%+ marginal farmers lack bank loan access due to zero KYC history\n• Trapped in illegal loan-shark debt (দাদন) with 80%+ interest rates\n• No digital mechanism to prove farm productivity to commercial banks.",
         Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6)),
    ]

    for title, desc, l, t, w, h in cards_data_s2:
        add_card(s2, l, t, w, h, title=title)
        tb = s2.shapes.add_textbox(l + Inches(0.2), t + Inches(0.6), w - Inches(0.4), h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 3: THE SOLUTION — 7-ROLE UNIFIED SUPER-APP
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, LIGHT_BG)
    add_header(s3, "The Solution: AgroLinkBD 7-Role Adaptive Ecosystem")

    # Left Column: The 7-Role Breakdown
    add_card(s3, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.4), title="Unified Adaptive Architecture for All 7 Stakeholders")
    roles_list = [
        ("1. Crop Farmer (ফসল চাষী):", "Land diary, plant doctor, fertilizer dosage, weather alerts & live bazaar."),
        ("2. Fish Farmer (মৎস্যচাষী):", "Multi-pond telemetry, Gemini AI fish doctor, FCR simulator, auction trade."),
        ("3. Wholesale Buyer (আড়তদার):", "Live auction bidding, corporate contract farming, bulk RFQ procurement."),
        ("4. Fleet Driver (পরিবহন চালক):", "Upazila-to-district truck load board, oxygenated live fish transit & escrow."),
        ("5. Machinery & Input Dealer:", "Combine harvester, tractor & drone rental sharing, QR seed store."),
        ("6. Hatchery Operator (হ্যাচারি):", "Fingerling quality certification, spawn calculator, broodstock tracking."),
        ("7. Agro-Doctor (বিশেষজ্ঞ):", "Live telemedicine video consultations, e-prescriptions, soil/water lab tests.")
    ]
    tb_roles = s3.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(6.1), Inches(4.6))
    tf_r = tb_roles.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    for idx, (r_title, r_desc) in enumerate(roles_list):
        p = tf_r.paragraphs[0] if idx == 0 else tf_r.add_paragraph()
        p.space_after = Pt(4)
        run_t = p.add_run()
        run_t.text = r_title + " "
        run_t.font.bold = True
        run_t.font.size = Pt(11)
        run_t.font.color.rgb = PRIMARY_EMERALD
        run_d = p.add_run()
        run_d.text = r_desc
        run_d.font.size = Pt(10.5)
        run_d.font.color.rgb = TEXT_DARK

    # Right Column: Key Architectural Highlights
    add_card(s3, Inches(7.6), Inches(1.5), Inches(4.9), Inches(5.4), title="Core System Pillars")
    pillars = [
        ("⚡ Flutter Clean Architecture", "High-performance reactive frontend supporting both English & Bengali (Hind Siliguri)."),
        ("🔥 Firebase Real-Time Cloud", "Instant real-time auction synchronization and live stock status with 400+ lines of RBAC rules."),
        ("🧬 Gemini 2.5 Flash Vision AI", "Multimodal pathology diagnostics processing image pixels directly for disease diagnosis."),
        ("💳 NFC & QR Hardware Integration", "Physical AgroCard bridging digital literacy barriers for unbanked marginal farmers."),
        ("🛡️ Escrow-Secured Payments", "SSLCommerz multi-channel gateway holding funds safely until buyer goods acceptance.")
    ]
    tb_pil = s3.shapes.add_textbox(Inches(7.8), Inches(2.1), Inches(4.5), Inches(4.6))
    tf_p = tb_pil.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = 0
    for idx, (p_title, p_desc) in enumerate(pillars):
        p = tf_p.paragraphs[0] if idx == 0 else tf_p.add_paragraph()
        p.space_after = Pt(6)
        run_t = p.add_run()
        run_t.text = p_title + "\n"
        run_t.font.bold = True
        run_t.font.size = Pt(11.5)
        run_t.font.color.rgb = DEEP_TEAL
        run_d = p.add_run()
        run_d.text = p_desc
        run_d.font.size = Pt(10.5)
        run_d.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 4: CORE INNOVATION #1 — GEMINI VISION AI DOCTOR
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, LIGHT_BG)
    add_header(s4, "Innovation #1: Ultra Pro Max AI Fish & Crop Doctor")

    ai_features = [
        ("🚫 Anti-Deception Guardrail", 
         "• Rejects human selfies, cars, animals, and blur photos.\n• Politely guides farmers with Bengali photography tips to capture clear disease lesions."),
        ("✅ 100% Healthy Specimen Vitality", 
         "• Recognizes disease-free specimens (clear eyes, intact fins, vibrant scales).\n• Issues health certificates & preventive feeding advice instead of prescribing unnecessary antibiotics."),
        ("📐 Decimal Chemical Dosage Engine", 
         "• Dynamically calculates exact dosages for Lime (চুন), Salt (লবণ), Zeolite (জিওলাইট), Potash, and Oxytetracycline.\n• Tailored to exact pond decimal area & water depth with zero overdosing risks."),
        ("🌱 Organic & Herbal Remedies", 
         "• Provides natural herbal treatments (Neem leaf, Turmeric, Mahogany oil extracts) alongside modern clinical prescriptions.")
    ]

    for idx, (title, desc) in enumerate(ai_features):
        col = idx % 2
        row = idx // 2
        l = Inches(0.8) if col == 0 else Inches(6.8)
        t = Inches(1.5) if row == 0 else Inches(4.3)
        w = Inches(5.6) if col == 0 else Inches(5.7)
        h = Inches(2.6)
        
        add_card(s4, l, t, w, h, title=title)
        tb = s4.shapes.add_textbox(l + Inches(0.2), t + Inches(0.6), w - Inches(0.4), h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 5: CORE INNOVATION #2 — LIVE AUCTIONS & CONTRACT FARMING
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, LIGHT_BG)
    add_header(s5, "Innovation #2: Live Wholesale Floor, Auctions & Contract Farming")

    bazaar_cards = [
        ("⚡ Live Digital Open Auctions (লাইভ নিলাম)", 
         "• Farmers broadcast bulk crop/fish harvests to verified national aratdars.\n• Transparent live bidding timer maximizes final sale price.\n• Direct buyer escrow deposit prevents payment defaults.",
         Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.6)),
        ("📜 Corporate Forward Contracts (চুক্তি চাষ)", 
         "• Formal pre-harvest buyback agreements with food MNCs (PRAN, ACI, Square).\n• Locked-in guaranteed pricing protects farmers from sudden seasonal price crashes.\n• Digital milestone tracking with advance token mobilization.",
         Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.6)),
        ("🛍️ 2-Tab Dual Trading Floor", 
         "• Tab 1: Live Wholesale Market Floor with dynamic search, commodity price ticker & 1-tap cart.\n• Tab 2: My Shop & Inventory Hub with live stock valuation & instant inline price/quantity editing.",
         Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6)),
        ("📋 Bulk Buyer RFQ Tendering Board", 
         "• Super-shops (Shwapno, Meena Bazar) and exporters post large procurement demands.\n• Farmers & cooperatives submit direct wholesale proposals without middlemen.",
         Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6))
    ]

    for title, desc, l, t, w, h in bazaar_cards:
        add_card(s5, l, t, w, h, title=title)
        tb = s5.shapes.add_textbox(l + Inches(0.2), t + Inches(0.6), w - Inches(0.4), h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 6: CORE INNOVATION #3 — 14-DAY AI PRICE RADAR & WEATHER
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, LIGHT_BG)
    add_header(s6, "Innovation #3: 14-Day AI Price Forecast Radar & Weather Telemetry")

    radar_cards = [
        ("📈 14-Day Machine Learning Price Radar",
         "• Predictive time-series model forecasting commodity price trajectories.\n• Covers 64 district wholesale markets (Potato, Rice, Onion, Chili, Mango, Fish).\n• Empowers farmers to time their harvest for peak market pricing windows.",
         Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.6)),
        ("📊 Real-Time Commodity Price Ticker",
         "• Scrolling live market rates with daily price trend indicators (+/- %).\n• Sourced from national trading hubs (Karwan Bazar, Bogura, Rajshahi, Chandpur).\n• Eradicates informational asymmetry exploited by local middlemen.",
         Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.6)),
        ("⛈️ Hyper-Local Disaster & Weather Alerts",
         "• Upazila-level real-time meteorological tracking.\n• Instant push warnings for Kalbaishakhi storms, flash floods, cold waves & droughts.\n• Actionable crop/pond protection protocols delivered in Bengali.",
         Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6)),
        ("🧪 Soil & Water Quality Telemetry",
         "• Soil NPK & pH lab test submission with crop suitability recommendations.\n• Aquaculture water quality tracking (pH, DO, Ammonia, Nitrite, Alkalinity).",
         Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6))
    ]

    for title, desc, l, t, w, h in radar_cards:
        add_card(s6, l, t, w, h, title=title)
        tb = s6.shapes.add_textbox(l + Inches(0.2), t + Inches(0.6), w - Inches(0.4), h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 7: CORE INNOVATION #4 — UPAZILA LOGISTICS & MECHANIZATION
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, LIGHT_BG)
    add_header(s7, "Innovation #4: Upazila Logistics & Farm Mechanization Fleet")

    log_cards = [
        ("🚚 On-Demand Upazila Truck Booking",
         "• Direct hire of 10-Ton, 5-Ton, 3-Ton trucks & pickup vans at upazila hubs.\n• Live GPS route tracking & digital challan generation.\n• Cuts perishable transit delays, reducing post-harvest waste by 30%.",
         Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.6)),
        ("🐟 Oxygenated Live Fish Transit Tanks",
         "• Specialized booking for trucks equipped with live aerated water tanks.\n• Eliminates transit fish mortality, enabling live fish sales at premium market rates in major metropolitan cities.",
         Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.6)),
        ("🚜 Farm Machinery & Drone Rental Sharing",
         "• Hourly/Bigha rental of Combine Harvesters, Tractors, Power Tillers & Seeders.\n• Agricultural Spraying Drones for rapid, uniform pesticide application.\n• Reduces manual labor costs & harvest time by 35%.",
         Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6)),
        ("🛡️ Escrow Deposit & Trip Insurance",
         "• Platform escrow protects both drivers and farmers against cargo damage or payment cancellation.\n• Integrated driver trip ledger & performance rating.",
         Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6))
    ]

    for title, desc, l, t, w, h in log_cards:
        add_card(s7, l, t, w, h, title=title)
        tb = s7.shapes.add_textbox(l + Inches(0.2), t + Inches(0.6), w - Inches(0.4), h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 8: CORE INNOVATION #5 — SMART AGROCARD & FINTECH
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, LIGHT_BG)
    add_header(s8, "Innovation #5: Hardware Smart AgroCard & Loan KYC Dossier")

    fin_cards = [
        ("💳 NFC & QR Hardware Smart AgroCard",
         "• Physical/digital identity card carrying a 6-digit PIN-secured digital wallet.\n• Solves digital literacy barriers: 1-tap offline identification at local dealers.\n• Multi-tier encryption protects farmer funds against unauthorized access.",
         Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.6)),
        ("🏦 Automated Bank Loan KYC File Generator",
         "• Converts farm harvest logs, land size, and transaction turnover into an official Bank Loan Project Report.\n• Accepted by Krishi Bank, Rajshahi Krishi Unnayan Bank, Sonali Bank, and BRAC Bank.\n• Eliminates collateral harassment and reliance on local loan sharks (দাদন).",
         Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.6)),
        ("📱 Integrated SSLCommerz Digital Payments",
         "• Multi-channel checkout supporting bKash, Nagad, Rocket, Upay, Cards, and Bank Wire.\n• Seamless wallet-to-wallet transfers with zero transaction friction.",
         Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6)),
        ("🩺 Telemedicine & Government Subsidy Radar",
         "• Live video consultations with verified aquaculture and agronomy scientists.\n• Real-time notifications of Ministry of Agriculture seed grants & solar subsidies.",
         Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6))
    ]

    for title, desc, l, t, w, h in fin_cards:
        add_card(s8, l, t, w, h, title=title)
        tb = s8.shapes.add_textbox(l + Inches(0.2), t + Inches(0.6), w - Inches(0.4), h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 9: BUSINESS MODEL & 6 MONETIZATION STREAMS
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, LIGHT_BG)
    add_header(s9, "Business Model: 6 Sustainable Commercial Revenue Streams")

    # Table of 6 Revenue Streams
    t_shape = s9.shapes.add_table(7, 4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.4))
    tbl = t_shape.table
    tbl.columns[0].width = Inches(3.2)
    tbl.columns[1].width = Inches(3.0)
    tbl.columns[2].width = Inches(3.2)
    tbl.columns[3].width = Inches(2.3)

    headers = ["Revenue Stream", "Target Segment", "Monetization Mechanism", "Yr 2 Est. Revenue"]
    for c_idx, h in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_EMERALD
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    b_rows = [
        ("1. B2B Wholesale Trade Commission", "Wholesalers / Aratdars", "1.5% – 2.5% transaction commission on auction/contract sales", "৳ 3.20 Crore"),
        ("2. Telemedicine Specialist Consultation", "Farmers seeking experts", "৳ 30 – ৳ 50 per live specialist video call token", "৳ 45.00 Lakh"),
        ("3. Logistics & Machinery Rental Escrow", "Fleet drivers & equipment owners", "3.0% – 5.0% platform safety and escrow fee", "৳ 85.00 Lakh"),
        ("4. B2B Corporate Sponsored Brand Deals", "Seed/Feed/Agro MNCs (ACI, Lal Teer)", "Sponsored brand placement & distributor zones", "৳ 60.00 Lakh"),
        ("5. VIP Pass & Intelligence Subscription", "Commercial fish & crop producers", "৳ 199/month for 14-day price radar & bank report export", "৳ 72.00 Lakh"),
        ("6. Microfinance Loan Origination Fee", "Partner Commercial Banks / MFIs", "0.5% – 1.0% origination fee on disbursed credit dossiers", "৳ 90.00 Lakh"),
    ]

    for r_idx, row in enumerate(b_rows):
        bg = RGBColor(232, 245, 233) if r_idx % 2 == 1 else WHITE
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_DARK
            if c_idx == 3:
                p.font.bold = True
                p.font.color.rgb = PRIMARY_EMERALD
                p.alignment = PP_ALIGN.RIGHT

    # Bottom Summary Card
    add_card(s9, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.8), bg_color=CARD_DARK_BG, border_color=PRIMARY_EMERALD)
    tb_sum = s9.shapes.add_textbox(Inches(1.0), Inches(5.25), Inches(11.3), Inches(1.5))
    tf_s = tb_sum.text_frame
    tf_s.word_wrap = True
    p = tf_s.paragraphs[0]
    p.text = "★ Total Projected Year 2 Revenue: ৳ 6.72 Crore | Net Profit Margin: 34.2%"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD_ACCENT

    p = tf_s.add_paragraph()
    p.text = "Sustainable Unit Economics: Zero marginal server cost with Google Cloud serverless architecture. High transactional volume from B2B wholesale trade provides immense organic cash flow without relying on donations or charity."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_LIGHT

    # =========================================================================
    # SLIDE 10: TECHNICAL ARCHITECTURE & SECURITY
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, LIGHT_BG)
    add_header(s10, "Technical System Architecture & Enterprise Security")

    tech_cards = [
        ("📱 Client Application Layer",
         "• Flutter Framework (Cross-Platform Android, iOS, Web)\n• Reactive GetX Controllers + Multi-Provider State Management\n• Optimized offline caching for rural low-bandwidth conditions\n• Dual Dark/Light Mode + Bengali Hind Siliguri Typography",
         Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.6)),
        ("🔥 Cloud Database & Backend",
         "• Google Firebase Firestore real-time document streaming\n• Serverless Cloud Functions handling transaction escrows\n• Firebase Storage with encrypted image chunk uploads\n• Auto-scaling capacity for 1,000,000+ simultaneous bids",
         Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.6)),
        ("🧬 Artificial Intelligence Engine",
         "• Google Gemini 2.5 Flash Multimodal Vision API\n• Real-time visual pathology analysis for crops & fish\n• Multi-model heuristic fallback guaranteeing 99.9% uptime\n• Time-series machine learning model for 14-day price forecasting",
         Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6)),
        ("🛡️ Enterprise Security & Fintech",
         "• 400+ lines of Granular Firestore Security Rules with RBAC\n• Digital Wallet PINs hashed with irreversible salt + AES-256\n• SSLCommerz PCI-DSS Level 1 certified gateway compliance\n• Hardware NFC Card UID token authentication",
         Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6))
    ]

    for title, desc, l, t, w, h in tech_cards:
        add_card(s10, l, t, w, h, title=title)
        tb = s10.shapes.add_textbox(l + Inches(0.2), t + Inches(0.6), w - Inches(0.4), h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 11: COMPETITIVE BENCHMARKING (WHY AGROLINKBD WINS)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11, LIGHT_BG)
    add_header(s11, "Competitive Advantage: Why AgroLinkBD Wins Nationwide")

    t_shape11 = s11.shapes.add_table(7, 5, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.4))
    tbl11 = t_shape11.table
    tbl11.columns[0].width = Inches(3.7)
    tbl11.columns[1].width = Inches(2.2)
    tbl11.columns[2].width = Inches(1.8)
    tbl11.columns[3].width = Inches(1.8)
    tbl11.columns[4].width = Inches(2.2)

    comp_headers = ["Key Feature / Capability", "AgroLinkBD", "Krishi Batayon", "iFarmer", "Chaldal / Shodagor"]
    for c_idx, h in enumerate(comp_headers):
        cell = tbl11.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_EMERALD if c_idx != 1 else GOLD_ACCENT
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    comp_rows = [
        ("1. Gemini Vision AI Doctor (Fish & Crop)", "✅ Multimodal + Dosage", "❌ Static Text Only", "❌ None", "❌ None"),
        ("2. Live Digital Auctions & Contract Farming", "✅ Real-Time Open Floor", "❌ None", "❌ None", "❌ Fixed Retail Only"),
        ("3. 14-Day AI Price Forecast Radar", "✅ 64-District ML Model", "❌ None", "❌ None", "❌ None"),
        ("4. Smart NFC AgroCard & Bank KYC Dossier", "✅ Hardware+App Dossier", "❌ None", "❌ Semi-Digital", "❌ None"),
        ("5. Upazila Logistics & Machinery Sharing", "✅ On-Demand GPS Fleet", "❌ None", "❌ None", "❌ City Logistics Only"),
        ("6. 7-Role Adaptive Integrated Super-App", "✅ Full Ecosystem", "❌ 1 Role Only", "❌ 1 Role Only", "❌ E-commerce Only"),
    ]

    for r_idx, row in enumerate(comp_rows):
        bg = RGBColor(232, 245, 233) if r_idx % 2 == 1 else WHITE
        for c_idx, val in enumerate(row):
            cell = tbl11.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_DARK
            if c_idx == 1:
                p.font.bold = True
                p.font.color.rgb = PRIMARY_EMERALD

    # =========================================================================
    # SLIDE 12: EXPECTED IMPACT & CALL TO ACTION (CLOSING SLIDE)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12, DARK_BG)

    # Header
    tbox = s12.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.2))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Transforming Bangladesh Agriculture: The Vision for 2041"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 4 Impact Pillars
    impacts = [
        ("+25% to 40%", "Farmer Profit Increase", "Eliminating multi-tier middlemen through transparent live auctions and corporate contracts."),
        ("-30%", "Loss & Spoilage Reduction", "Rapid Gemini Vision pathology diagnosis and cold-chain upazila transport."),
        ("100,000+", "Farmers Bank Inclusion", "NFC Smart AgroCards providing automated KYC credit dossiers for collateral-free bank loans."),
        ("SDG 1, 2, 8, 9", "National Development", "Directly driving Smart Bangladesh 2041, Zero Hunger, No Poverty, and Rural Economic Growth.")
    ]

    for idx, (stat, title, desc) in enumerate(impacts):
        l = Inches(0.8 + idx * 2.95)
        t = Inches(2.2)
        w = Inches(2.8)
        h = Inches(3.2)
        add_card(s12, l, t, w, h, bg_color=CARD_DARK_BG, border_color=PRIMARY_EMERALD)
        
        tb = s12.shapes.add_textbox(l + Inches(0.15), t + Inches(0.2), w - Inches(0.3), h - Inches(0.4))
        tf_i = tb.text_frame
        tf_i.word_wrap = True
        p = tf_i.paragraphs[0]
        p.text = stat
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GOLD_ACCENT
        p.alignment = PP_ALIGN.CENTER
        
        p = tf_i.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.CENTER
        
        p = tf_i.add_paragraph()
        p.text = desc
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

    # Final Closing Punchline Box
    fin_box = s12.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2))
    tf_f = fin_box.text_frame
    tf_f.word_wrap = True
    p = tf_f.paragraphs[0]
    p.text = "🌾 AgroLinkBD is not just an application — it is the digital infrastructure for Bangladesh's Agricultural Revolution."
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(129, 199, 132)
    p.alignment = PP_ALIGN.CENTER

    p = tf_f.add_paragraph()
    p.text = "Thank You | Questions & Live Interactive Demonstration"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Save Presentation
    prs.save(output_filename)
    print("SUCCESS: Created PowerPoint Presentation at " + output_filename)

if __name__ == "__main__":
    create_presentation_deck(r"d:\App\AgroLinkBD\AgroLinkBD_Innovation_for_Bangladesh_Presentation.pptx")
